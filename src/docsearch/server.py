import logging
import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree

from typing import Annotated
from mcp.server.fastmcp import FastMCP
from pydantic import Field

logging.basicConfig(stream=sys.stderr, level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger("docsearch")

mcp = FastMCP("docsearch")

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".odt", ".ods", ".odp", ".rtf", ".epub"}
MAX_OUTPUT_CHARS = 40_000


def extract_pdf(path: Path) -> list[tuple[str, str]]:
    import pymupdf
    sections = []
    with pymupdf.open(str(path)) as doc:
        for i, page in enumerate(doc, 1):
            text = page.get_text("text").strip()
            if text:
                sections.append((f"page {i}", text))
    return sections


def extract_docx(path: Path) -> list[tuple[str, str]]:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cells)
            if line.strip():
                parts.append(line)
    if parts:
        return [("document", "\n".join(parts))]
    return []


def extract_pptx(path: Path) -> list[tuple[str, str]]:
    from pptx import Presentation
    sections = []
    prs = Presentation(str(path))
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    line = "\t".join(cells)
                    if line.strip():
                        parts.append(line)
        if parts:
            sections.append((f"slide {i}", "\n".join(parts)))
    return sections


def extract_xlsx(path: Path) -> list[tuple[str, str]]:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    sections = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = "\t".join(cells)
            if line.strip():
                lines.append(line)
        if lines:
            sections.append((f"sheet '{sheet_name}'", "\n".join(lines)))
    wb.close()
    return sections


def extract_odt(path: Path) -> list[tuple[str, str]]:
    from odf.opendocument import load
    from odf import text as odf_text, teletype
    doc = load(str(path))
    parts = [teletype.extractText(p) for p in doc.getElementsByType(odf_text.P) if teletype.extractText(p).strip()]
    if parts:
        return [("document", "\n".join(parts))]
    return []


def extract_ods(path: Path) -> list[tuple[str, str]]:
    from odf.opendocument import load
    from odf import table as odf_table, teletype, text as odf_text
    doc = load(str(path))
    sections = []
    for sheet in doc.getElementsByType(odf_table.Table):
        sheet_name = sheet.getAttribute("name")
        lines = []
        for row in sheet.getElementsByType(odf_table.TableRow):
            cells = []
            for cell in row.getElementsByType(odf_table.TableCell):
                repeat = int(cell.getAttribute("numbercolumnsrepeated") or 1)
                value = teletype.extractText(cell).strip()
                cells.extend([value] * repeat)
            line = "\t".join(cells).rstrip("\t")
            if line.strip():
                lines.append(line)
        if lines:
            sections.append((f"sheet '{sheet_name}'", "\n".join(lines)))
    return sections


def extract_odp(path: Path) -> list[tuple[str, str]]:
    from odf.opendocument import load
    from odf import draw, teletype, text as odf_text
    doc = load(str(path))
    sections = []
    for i, page in enumerate(doc.getElementsByType(draw.Page), 1):
        parts = [teletype.extractText(p) for p in page.getElementsByType(odf_text.P) if teletype.extractText(p).strip()]
        if parts:
            sections.append((f"slide {i}", "\n".join(parts)))
    return sections


def extract_rtf(path: Path) -> list[tuple[str, str]]:
    from striprtf.striprtf import rtf_to_text
    raw = path.read_bytes().decode("utf-8", errors="replace")
    text = "\n".join(line for line in rtf_to_text(raw).splitlines() if line.strip())
    if text:
        return [("document", text)]
    return []


def extract_epub(path: Path) -> list[tuple[str, str]]:
    tag_re = re.compile(r"<[^>]+>")
    sections = []
    with zipfile.ZipFile(str(path), "r") as zf:
        container = ElementTree.fromstring(zf.read("META-INF/container.xml"))
        ns = {"c": "urn:oasis:names:tc:opendocument:xmlns:container"}
        opf_path = container.find(".//c:rootfile", ns).get("full-path")
        opf_dir = opf_path.rsplit("/", 1)[0] + "/" if "/" in opf_path else ""
        opf = ElementTree.fromstring(zf.read(opf_path))
        opf_ns = {"opf": "http://www.idpf.org/2007/opf"}
        manifest = {item.get("id"): item.get("href") for item in opf.findall(".//opf:manifest/opf:item", opf_ns)}
        spine_ids = [ref.get("idref") for ref in opf.findall(".//opf:spine/opf:itemref", opf_ns)]
        for i, idref in enumerate(spine_ids, 1):
            href = manifest.get(idref)
            if not href:
                continue
            full_path = opf_dir + href if not href.startswith("/") else href.lstrip("/")
            try:
                html = zf.read(full_path).decode("utf-8", errors="replace")
            except KeyError:
                continue
            text = tag_re.sub("", html)
            lines = [l.strip() for l in text.splitlines() if l.strip()]
            if lines:
                sections.append((f"chapter {i}", "\n".join(lines)))
    return sections


def extract_text(path: Path) -> list[tuple[str, str]]:
    extractors = {
        ".pdf": extract_pdf, ".docx": extract_docx, ".pptx": extract_pptx, ".xlsx": extract_xlsx,
        ".odt": extract_odt, ".ods": extract_ods, ".odp": extract_odp, ".rtf": extract_rtf, ".epub": extract_epub,
    }
    ext = path.suffix.lower()
    if ext not in extractors:
        raise ValueError(f"Unsupported file type: {ext}")
    return extractors[ext](path)


def parse_numeric_range(s: str) -> set[int]:
    nums = set()
    for part in s.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            nums.update(range(int(start), int(end) + 1))
        else:
            nums.add(int(part))
    return nums


def extract_section_number(label: str) -> int | None:
    m = re.match(r".*?(\d+)$", label)
    return int(m.group(1)) if m else None


def filter_sections(sections: list[tuple[str, str]], range_str: str, ext: str) -> tuple[str | None, list[tuple[str, str]]]:
    if ext in (".docx", ".odt", ".rtf"):
        line_range = parse_numeric_range(range_str)
        filtered = []
        for label, text in sections:
            lines = text.split("\n")
            selected = [lines[j - 1] for j in sorted(line_range) if 1 <= j <= len(lines)]
            if selected:
                lo, hi = min(line_range), max(line_range)
                filtered.append((f"{label} lines {lo}-{hi}", "\n".join(selected)))
        return None, filtered

    if ext in (".xlsx", ".ods"):
        names = set()
        indices = set()
        row_range = None
        for part in range_str.split(","):
            part = part.strip().strip("'\"")
            if ":" in part:
                sheet_part, row_part = part.split(":", 1)
                row_range = parse_numeric_range(row_part)
                part = sheet_part
            try:
                indices.add(int(part))
            except ValueError:
                if "-" not in part:
                    names.add(part)
                else:
                    try:
                        indices |= parse_numeric_range(part)
                    except ValueError:
                        names.add(part)
        filtered = []
        for i, (label, text) in enumerate(sections, 1):
            sheet_name = label.removeprefix("sheet '").removesuffix("'")
            if i in indices or sheet_name in names:
                if row_range is not None:
                    lines = text.split("\n")
                    selected = [lines[j - 1] for j in sorted(row_range) if 1 <= j <= len(lines)]
                    lo, hi = min(row_range), max(row_range)
                    filtered.append((f"{label} rows {lo}-{hi}", "\n".join(selected)))
                else:
                    filtered.append((label, text))
        return None, filtered

    indices = parse_numeric_range(range_str)
    return None, [(l, t) for l, t in sections if extract_section_number(l) in indices]


@mcp.tool()
def docgrep(
    directory: Annotated[str, Field(description="Path to the directory to search. Searched recursively, skipping hidden directories.")],
    pattern: Annotated[str, Field(description="Python regex pattern to match against line content.")],
    case_sensitive: Annotated[bool, Field(description="Use case-sensitive matching. Default is case-insensitive.")] = False,
    file_types: Annotated[list[str] | None, Field(description="Limit to specific file types, with or without dots (e.g. ['pdf', 'docx']). Defaults to all supported types.")] = None,
    max_results: Annotated[int, Field(description="Maximum number of matching lines to return.")] = 100,
) -> str:
    """Search through document files (PDF, DOCX, PPTX, XLSX, ODT, ODS, ODP, RTF, EPUB) for text matching a regex pattern. Returns grep-like output: filepath:section:matching_line"""
    root = Path(directory).resolve()
    if not root.is_dir():
        raise ValueError(f"Not a directory: {directory}")

    flags = 0 if case_sensitive else re.IGNORECASE
    try:
        regex = re.compile(pattern, flags)
    except re.error as e:
        raise ValueError(f"Invalid regex pattern: {e}")

    allowed_exts = SUPPORTED_EXTENSIONS
    if file_types:
        allowed_exts = {ext if ext.startswith(".") else f".{ext}" for ext in file_types} & SUPPORTED_EXTENSIONS

    results = []
    for filepath in sorted(root.rglob("*")):
        if len(results) >= max_results:
            break
        if not filepath.is_file():
            continue
        if filepath.suffix.lower() not in allowed_exts:
            continue
        if any(part.startswith(".") for part in filepath.relative_to(root).parts):
            continue

        try:
            sections = extract_text(filepath)
        except Exception as e:
            logger.error(f"Error extracting {filepath}: {e}")
            results.append(f"{filepath.relative_to(root)}:error:{e}")
            continue

        for section_label, text in sections:
            if len(results) >= max_results:
                break
            for line in text.splitlines():
                if len(results) >= max_results:
                    break
                if regex.search(line):
                    rel = filepath.relative_to(root)
                    results.append(f"{rel}:{section_label}:{line.strip()}")

    if not results:
        return "No matches found."
    return "\n".join(results)


@mcp.tool()
def docread(
    filepath: Annotated[str, Field(description="Path to the document file.")],
    range: Annotated[str | None, Field(
        description="Filter output to specific sections. Format depends on file type:\n"
        "- PDF: page numbers (e.g. '1-3', '1,3,5-7')\n"
        "- PPTX/ODP: slide numbers (e.g. '2-4')\n"
        "- EPUB: chapter numbers (e.g. '1-5')\n"
        "- XLSX/ODS: sheet name or 1-based index, with optional row range after colon (e.g. '1', 'Sheet1', '1:1-100', 'Revenue:50-200')\n"
        "- DOCX/ODT/RTF: line numbers (e.g. '1-50', '100-200')"
    )] = None,
) -> str:
    """Read and extract text from a document file (PDF, DOCX, PPTX, XLSX, ODT, ODS, ODP, RTF, EPUB). Returns sections formatted as '=== section_label ===\\ntext'. Output auto-truncated at 40000 chars."""
    path = Path(filepath).resolve()
    if not path.is_file():
        raise ValueError(f"File not found: {filepath}")
    if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {path.suffix}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")

    sections = extract_text(path)

    warning = None
    if range:
        warning, sections = filter_sections(sections, range, path.suffix.lower())

    if not sections:
        return "No text content found."

    output = []
    if warning:
        output.append(warning + "\n")
    for label, text in sections:
        output.append(f"=== {label} ===\n{text}\n")
    result = "\n".join(output)
    if len(result) > MAX_OUTPUT_CHARS:
        result = result[:MAX_OUTPUT_CHARS] + f"\n\n... output truncated at {MAX_OUTPUT_CHARS} chars. Use range to narrow results."
    return result


def main():
    mcp.run(transport="stdio")
